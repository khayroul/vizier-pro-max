"""Assemble a PPTX deck from structured slide content."""
from __future__ import annotations

from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)


def run(
    *,
    title: str,
    slides: list[dict[str, object]],
    output_path: str,
) -> dict[str, str]:
    """Create a PPTX presentation from structured slide data."""
    if not title:
        msg = "title is required"
        raise ValueError(msg)
    if not output_path:
        msg = "output_path is required"
        raise ValueError(msg)

    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    deck = Presentation()
    deck.core_properties.title = title

    for slide_data in slides:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = str(slide_data.get("title", ""))

        try:
            slide.placeholders[1].text = str(slide_data.get("body", ""))
        except KeyError:
            pass

        image_path = str(slide_data.get("image_path", ""))
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(image_path, Inches(5), Inches(1.5), width=Inches(4))

        notes = str(slide_data.get("notes", ""))
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(output))
    logger.info("Assembled PPTX", output_path=str(output), slide_count=len(slides))
    return {"file_path": str(output)}
